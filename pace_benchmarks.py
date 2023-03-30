# -*- coding: utf-8 -*-

import os
import json
import glob
import sys
import pandas as pd
from pandas.plotting import table
from monty.json import MontyDecoder, MontyEncoder
from sklearn.metrics import mean_absolute_error, mean_squared_error
import numpy as np
from tqdm import tqdm
from matplotlib import pyplot as plt
import scienceplots
from matplotlib.ticker import MultipleLocator
from scipy.optimize import curve_fit
from scipy import interpolate
from phonopy import Phonopy
from phonopy.structure.atoms import PhonopyAtoms
from pymatgen.io.phonopy import get_pmg_structure, get_phonopy_structure
from phonopy.phonon.band_structure import get_band_qpoints_and_path_connections
from uberplot import (
    load_record,
    tabulize_record,
    in_plane_vectors,
    transform_energy_to_surface_energy,
    uber_fit,
    uber,
)
from pptx import Presentation
from pptx.util import Inches, Pt

import json
from monty.json import MontyDecoder, MontyEncoder
from sklearn.metrics import mean_absolute_error, mean_squared_error
import warnings

warnings.filterwarnings("ignore")

######setup plotting#####
colors = [
    "black",
    "dimgray",
    "lightcoral",
    "brown",
    "red",
    "orangered",
    "chocolate",
    "burlywood",
    "orange",
    "darkgoldenrod",
    "gold",
    "olive",
    "darkolivegreen",
    "chartreuse",
    "green",
    "turquoise",
    "darkcyan",
    "deepskyblue",
    "dodgerblue",
    "midnightblue",
    "blue",
    "blueviolet",
    "violet",
    "deeppink",
]
index = np.array(
    [
        0,
        4,
        8,
        12,
        16,
        20,
        1,
        5,
        9,
        13,
        17,
        21,
        2,
        6,
        10,
        14,
        18,
        22,
        3,
        7,
        11,
        15,
        19,
        23,
    ]
)
colors_ = []
for i in index:
    colors_.append(colors[i])
colors = colors_
markers = [
    ".",
    "o",
    "v",
    "^",
    "<",
    ">",
    "s",
    "*",
    "D",
    "1",
    "+",
    "x",
    ".",
    "o",
    "v",
    "^",
    "<",
    ">",
    "s",
    "*",
    "D",
    "1",
    "+",
    "x",
]


##### structure renaming ######
renaming = {
    "mp-753POSCAR": "hcp",
    "mp-23313POSCAR": "fcc",
    "mp-10021POSCAR": r"$\rho_1$",
    "mp-142POSCAR": r"$\rho_2$",
    "mp-29647POSCAR": r"$\rho_3$",
    "mp-568087POSCAR": r"$\rho_4$",
    "mp-568714POSCAR": r"$\rho_5$",
    "mp-605790POSCAR": r"$\rho_6$",
    "mp-639736POSCAR": r"$\rho_7$",
    "mp-640118POSCAR": r"$\rho_8$",
    "mp-84POSCAR": r"$\rho_9$",
    "mp-865373POSCAR": r"$\rho_{10}$",
    "crsi2": r"$CrSi_2$",
    "lisi": r"$LiSi$",
    "tisi2": r"$TiSi_2$",
    "la2sb": r"$La_2Sb$",
    "ptte": r"$PtTe$",
    "hg2pt": r"$Hg_2Pt$",
    "bapb3": r"$BaPb_3$",
    "hg2pt": r"$Hg_2Pt$",
}

#####save a picture of a table#####
def savetable(data: dict, name: str, dpi: int):
    """
    input a dictionary of a table; the name and dpi of output picture
    output a picture of the table
    """
    df = pd.DataFrame(data=data)
    fig = plt.figure(figsize=(5, 6))
    ax = fig.add_subplot(111, frame_on=False)
    ax.xaxis.set_visible(False)
    ax.yaxis.set_visible(False)
    table(ax, df, loc="center")
    plt.savefig(name, dpi=dpi)


#####get id for prototypes in GB,surface calculations#####
def get_ref_id(data_collection: list) -> dict:
    """
    get the ids for prototypes bcc, fcc, hcp, omega
    """
    ref_id = {}
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "icsd"
            and data_collection[i]["calc"] == "final"
            and "proto" in data_collection[i]["metadata"]
        ):
            if data_collection[i]["metadata"]["proto"] in [
                "bcc",
                "fcc",
                "hcp",
                "omega",
            ]:
                ref_id[data_collection[i]["metadata"]["proto"]] = i
    return ref_id


#####plot surface heatmap#####
def plot_surface_heatmap(ax, X, Y, Z):
    pc = ax.pcolormesh(X, Y, Z)
    ax.set_aspect("equal")
    im_ratio = (max(Y.ravel()) - min(Y.ravel())) / (max(X.ravel()) - min(X.ravel()))
    cbar = plt.colorbar(pc, ax=ax, fraction=0.047 * im_ratio, pad=0.04)
    return ax, cbar


####### Vacancy formation energy #######
########################################
def vacancy_formation(data_collection: list):
    """
    Vacancy formation energy analysis
    input data_collection::list
    output a picture of the table of the vacancy formation energies
    """
    vacancy_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "vacancies"
            and data_collection[i]["calc"] == "final"
            and not data_collection[i]["metadata"]["proto"] == "casi-alth"
            and not data_collection[i]["metadata"]["proto"] == "mp-865373POSCAR"
        ):
            vacancy_id.append(i)
    vacancy_ideal_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "vacancies_ideal"
            and data_collection[i]["calc"] == "final"
        ):
            vacancy_ideal_id.append(i)
    protos_withid = []
    protos_ = []
    for id in vacancy_id:
        protos_withid.append(
            data_collection[id]["metadata"]["proto"]
            + "_"
            + str(data_collection[id]["metadata"]["deleted_site"])
            + "|"
            + str(id)
        )
        protos_.append(
            data_collection[id]["metadata"]["proto"]
            + "_"
            + str(data_collection[id]["metadata"]["deleted_site"])
        )
    protos_ideal = {}
    for id in vacancy_ideal_id:
        protos_ideal[data_collection[id]["metadata"]["proto"]] = id
    deltaEs = []
    deltaEs_pred = []
    for proto in list(protos_withid):
        delid = int(proto.split("|")[1])
        idealid = protos_ideal[proto.split("_")[0]]
        N_proto = len(data_collection[idealid]["structure"].species)
        Eideal = data_collection[idealid]["energy"]
        Evac = data_collection[delid]["energy"]
        Eideal_pred = data_collection[idealid]["pace"]["energy"]
        Evac_pred = data_collection[delid]["pace"]["energy"]
        deltaE = Evac - (N_proto - 1) * Eideal / N_proto
        deltaE_pred = Evac_pred - (N_proto - 1) * Eideal_pred / N_proto
        deltaEs.append(deltaE)
        deltaEs_pred.append(deltaE_pred)

    # d = {
    #     "Structure_site": protos_,
    #     r"$E^{DFT} (eV)$": np.round(np.array(deltaEs), 5),
    #     r"$E^{pred} (eV)$": np.round(np.array(deltaEs_pred), 5),
    # }

    deltaEs = np.round(np.array(deltaEs), 3)
    deltaEs_pred = np.round(np.array(deltaEs_pred), 3)

    protos_collect_ = {}
    for i in np.arange(len(protos_)):
        name = protos_[i].split("_")[0]
        if name in renaming:
            name = renaming[name]
        if name not in protos_collect_.keys():
            protos_collect_[name] = [[], []]
        protos_collect_[name][0].append(deltaEs[i])
        protos_collect_[name][1].append(deltaEs_pred[i])

    protos_collect = []
    deltaEs_collect = []
    deltaEs_pred_collect = []
    discrepancies_collect = []

    for key, value in protos_collect_.items():
        protos_collect.append(key)
        deltaEs_collect.append(", ".join([str(f) for f in value[0]]))
        deltaEs_pred_collect.append(", ".join([str(f) for f in value[1]]))
        dis = (np.array(value[1]) - np.array(value[0])) / np.array(value[0]) > 0.15
        mapping = {False: "", True: "+"}
        dis = [mapping[item] for item in dis]
        dis = ",".join(dis)
        if dis == ",":
            dis = ""
        discrepancies_collect.append(dis)
    d = {
        "proto": protos_collect,
        r"$E^{DFT} (eV)$": deltaEs_collect,
        r"$E^{pred} (eV)$": deltaEs_pred_collect,
        "discrepancy": discrepancies_collect,
    }

    savetable(d, "Vacancy_formation_pace.png", dpi=200)
    return


####### Volumetric deformation #######
#############################
def volumetric_deformation(
    data_collection: list, ref_omega_dft: float, ref_omega_pace: float
):
    """
    Volumetric deformation energy analysis
    input data_collection::list
    ref_omega_dft::float   reference energy (per atom) for omega structure from dft calculation
    ref_omega_pace::float  reference energy (per atom) for omega structure from pace potential
    output pictures of Energy curves with deformation
    """
    volumetric_id = []
    for i in np.arange(len(data_collection)):
        if data_collection[i]["metadata"]["perturbation"] == "volumetric":
            volumetric_id.append(i)
    protos = set()
    for i in volumetric_id:
        protos.add(data_collection[i]["metadata"]["proto"])
    with plt.style.context("science"):
        plt.figure()
        for i in np.arange(len(protos)):
            Vs_ = []
            Es_ = []
            Epreds_ = []
            for j in volumetric_id:
                if data_collection[j]["metadata"]["proto"] == list(protos)[i]:
                    Vs_.append(
                        data_collection[j]["metadata"]["volume"]
                        / len(data_collection[j]["structure"].species)
                    )
                    Es_.append(
                        data_collection[j]["energy"]
                        / len(data_collection[j]["structure"].species)
                    )
                    Epreds_.append(
                        data_collection[j]["pace"]["energy"]
                        / len(data_collection[j]["structure"].species)
                    )
            Vs_ = np.array(Vs_)
            Vs = []
            Es = []
            Epreds = []
            for checkindex in np.arange(len(Vs_)):
                samemag = np.where(
                    abs(Vs_ - Vs_[checkindex]) < 5e-2
                )[0]
                if len(samemag) == 1:
                    Vs.append(Vs_[checkindex])
                    Es.append(Es_[checkindex])
                    Epreds.append(Epreds_[checkindex])
                else:
                    flag = 1
                    for mag in samemag:
                        if Es_[checkindex] > Es_[mag]:
                            flag = 0
                    if flag == 1:
                        Vs.append(Vs_[checkindex])
                        Es.append(Es_[checkindex])
                        Epreds.append(Epreds_[checkindex])
            Vs = np.array(Vs)
            Es = np.array(Es) - ref_omega_dft
            Epreds = np.array(Epreds) - ref_omega_pace
            order = np.argsort(Vs)
            label = list(protos)[i]
            if label in renaming:
                label = renaming[label]
            plt.scatter(
                Vs[order],
                Es[order],
                c=colors[i],
                marker=markers[i],
                s=0.2,
                label=label,
            )
            plt.plot(Vs[order], Epreds[order], c=colors[i])
        plt.xlabel(r"$\overline{V} \enspace (Å^{3}/atom)$")
        plt.ylabel(r"$\Delta E (eV/atom)$")
        plt.legend(bbox_to_anchor=(1.1, 1), ncol=2)
        plt.savefig("Volumetric_deformation_pace.png", dpi=200)
        plt.close()
    for iter in np.arange(4):
        with plt.style.context("science"):
            plt.figure()
            for i in np.arange(6):
                Vs = []
                Es = []
                Epreds = []
                for j in volumetric_id:
                    if (
                        data_collection[j]["metadata"]["proto"]
                        == list(protos)[i + iter * 6]
                    ):
                        Vs.append(
                            data_collection[j]["metadata"]["volume"]
                            / len(data_collection[j]["structure"].species)
                        )
                        Es.append(
                            data_collection[j]["energy"]
                            / len(data_collection[j]["structure"].species)
                        )
                        Epreds.append(
                            data_collection[j]["pace"]["energy"]
                            / len(data_collection[j]["structure"].species)
                        )
                Vs = np.array(Vs)
                Es = np.array(Es) - ref_omega_dft
                Epreds = np.array(Epreds) - ref_omega_pace
                order = np.argsort(Vs)
                label = list(protos)[i + iter * 6]
                if label in renaming:
                    label = renaming[label]
                plt.scatter(
                    Vs[order],
                    Es[order],
                    c=colors[i],
                    marker=markers[i],
                    s=3,
                    label=label,
                )
                plt.plot(Vs[order], Epreds[order], c=colors[i])
            plt.xlabel(r"$\overline{V} \enspace (Å^{3}/atom)$")
            plt.ylabel(r"$\Delta E \enspace (eV/atom)$")
            plt.legend(bbox_to_anchor=(0.5, -0.15))
            plt.savefig("Volumetric_deformation_" + str(iter) + "_pace.png", dpi=250)
            plt.close()


####### Grain boudary #######
#############################
def grain_boundary(data_collection: list):
    """
    Grain boundary energy analysis
    input data_collection::list
    output a picture of the table of the grain boundary energies
    """
    protoids = get_ref_id(data_collection)
    gb_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "gb"
            and data_collection[i]["calc"] == "final"
        ):
            gb_id.append(i)
    icsd_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "icsd"
            and data_collection[i]["calc"] == "final"
        ):
            icsd_id.append(i)
    deltaEs_gb = []
    deltaEspred_gb = []
    protos = []
    discrepancy = []
    for id in gb_id:
        protoid = 0
        if data_collection[id]["metadata"]["ref_elem"] == "Mo":
            protoid = protoids["bcc"]
            protos.append("bcc")
        if data_collection[id]["metadata"]["ref_elem"] == "Ti":
            protoid = protoids["hcp"]
            protos.append("hcp")
        N_gb = len(data_collection[id]["structure"].species)
        M_icsd = len(data_collection[protoid]["structure"].species)
        E_gb = data_collection[id]["energy"]
        E_gb_pred = data_collection[id]["pace"]["energy"]
        E_icsd = data_collection[protoid]["energy"]
        E_icsd_pred = data_collection[protoid]["pace"]["energy"]
        lattice = data_collection[id]["structure"].lattice.matrix
        A = np.linalg.norm(np.cross(lattice[0, :], lattice[1, :]))
        deltaE_gb = (E_gb - E_icsd * N_gb / M_icsd) / 2 / A
        deltaEpred_gb = (E_gb_pred - E_icsd_pred * N_gb / M_icsd) / 2 / A
        deltaEs_gb.append(deltaE_gb)
        deltaEspred_gb.append(deltaEpred_gb)
        if abs((deltaEpred_gb - deltaE_gb) / deltaE_gb) > 0.15:
            discrepancy.append("+")
        else:
            discrepancy.append("")
    planes = [data_collection[gb_id[i]]["metadata"]["plane"] for i in np.arange(8)]
    sigmas = [data_collection[gb_id[i]]["metadata"]["sigma"] for i in np.arange(8)]
    d = {
        "proto": protos,
        "plane": planes,
        "sigma": sigmas,
        r"$E^{DFT}(eV/Å^{2})$": np.round(np.array(deltaEs_gb), 3),
        r"$E^{pred}(eV/Å^{2})$": np.round(np.array(deltaEspred_gb), 3),
        "discrepancy": discrepancy,
    }

    savetable(d, "Grain_boundary_pace.png", 200)


####### Surface #######
#######################
def surface(data_collection: list):
    """
    Surface boundary energy analysis
    input data_collection::list
    output a picture of the table of the surface energies
    """
    surface_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "surfaces"
            and data_collection[i]["calc"] == "final"
        ):
            surface_id.append(i)
    icsd_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "icsd"
            and data_collection[i]["calc"] == "final"
        ):  # and data_collection[i]['metadata']['proto']=='bcc' or data_collection[i]['metadata']['proto']=='hcp':
            icsd_id.append(i)
    protoids = get_ref_id(data_collection)
    deltaEs_sf = []
    deltaEspred_sf = []
    discrepancy = []
    for id in surface_id:
        protoid = protoids[data_collection[id]["metadata"]["proto"]]
        N_sf = len(data_collection[id]["structure"].species)
        M_icsd = len(data_collection[protoid]["structure"].species)
        E_sf = data_collection[id]["energy"]
        E_sf_pred = data_collection[id]["pace"]["energy"]
        E_icsd = data_collection[protoid]["energy"]
        E_icsd_pred = data_collection[protoid]["pace"]["energy"]
        lattice = data_collection[id]["structure"].lattice.matrix
        A = np.linalg.norm(np.cross(lattice[0, :], lattice[1, :]))
        deltaE_sf = (E_sf - E_icsd * N_sf / M_icsd) / 2 / A
        deltaEpred_sf = (E_sf_pred - E_icsd_pred * N_sf / M_icsd) / 2 / A
        deltaEs_sf.append(deltaE_sf)
        deltaEspred_sf.append(deltaEpred_sf)
        if abs((deltaEpred_sf - deltaE_sf) / deltaE_sf) > 0.15:
            discrepancy.append("+")
        else:
            discrepancy.append("")
    protos = [
        data_collection[id]["metadata"]["proto"].split(".")[0] for id in surface_id
    ]
    planes = [data_collection[id]["metadata"]["miller"] for id in surface_id]
    d = {
        "proto": protos,
        "plane": planes,
        r"$E^{DFT}(eV/Å^{2})$": np.round(np.array(deltaEs_sf), 5),
        r"$E^{pred}(eV/Å^{2})$": np.round(np.array(deltaEspred_sf), 5),
        "discrepancy": discrepancy,
    }
    savetable(d, "Surfaces_pace.png", 200)


####### Strain #######
######################
def strain(data_collection: list, ref_omega_dft: float, ref_omega_pace):
    """
    Strain-Energy relations
    input
        data_collection::list
        ref_omega_dft::float   reference energy (per atom) for omega structure from dft calculation
        ref_omega_pace::float  reference energy (per atom) for omega structure from pace potential
    output pictures of the Energy-Strain curves (a picture contains different strain directions for the same prototype)
    """
    strain_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "strain"
            and data_collection[i]["calc"] == "final"
        ):
            strain_id.append(i)

    protos = set()
    for id in strain_id:
        protos.add(data_collection[id]["metadata"]["proto"])
    protos = list(protos)
    for proto in protos:
        strain_id = []
        for i in np.arange(len(data_collection)):
            if (
                data_collection[i]["metadata"]["perturbation"] == "strain"
                and data_collection[i]["calc"] == "final"
                and data_collection[i]["metadata"]["proto"] == proto
            ):
                strain_id.append(i)
        Es = np.zeros((len(strain_id), 6))
        i = 0
        for id in strain_id:
            Es[i] = np.array(data_collection[id]["metadata"]["special_direction"])
            i += 1
        groups = []
        already = np.zeros(len(Es))
        isidentical = abs(np.dot(Es, Es.T) - 1) < 1e-12
        for i in np.arange(len(Es)):
            if not already[i] == 1:
                thisgroup = np.where(isidentical[i] == True)
                groups.append(thisgroup[0].tolist())
                already[thisgroup] = 1
        i = 0
        with plt.style.context("science"):
            plt.figure()
            for group in groups:
                magnitudes_ = []
                Energies_ = []
                Energies_pred_ = []
                for groupid in group:
                    Natom = len(
                        data_collection[strain_id[groupid]]["structure"].species
                    )
                    magnitudes_.append(
                        data_collection[strain_id[groupid]]["metadata"]["magnitude"]
                    )
                    Energies_.append(
                        data_collection[strain_id[groupid]]["energy"] / Natom
                    )
                    Energies_pred_.append(
                        data_collection[strain_id[groupid]]["pace"]["energy"] / Natom
                    )
                magnitudes = []
                Energies = []
                Energies_pred = []
                magnitudes_ = np.array(magnitudes_)
                for checkindex in np.arange(len(magnitudes_)):
                    samemag = np.where(
                        abs(magnitudes_ - magnitudes_[checkindex]) < 1e-2
                    )[0]
                    if len(samemag) == 1:
                        magnitudes.append(magnitudes_[checkindex])
                        Energies.append(Energies_[checkindex])
                        Energies_pred.append(Energies_pred_[checkindex])
                    else:
                        flag = 1
                        for mag in samemag:
                            if Energies_[checkindex] > Energies_[mag]:
                                flag = 0
                        if flag == 1:
                            magnitudes.append(magnitudes_[checkindex])
                            Energies.append(Energies_[checkindex])
                            Energies_pred.append(Energies_pred_[checkindex])
                magnitudes = np.array(magnitudes)
                Energies = np.array(Energies)
                Energies_pred = np.array(Energies_pred)
                order = magnitudes.argsort()
                plt.scatter(
                    magnitudes[order],
                    Energies[order] - ref_omega_dft,
                    c=colors[i],
                    marker=markers[i],
                    s=2,
                    label=np.round(
                        data_collection[strain_id[group[0]]]["metadata"][
                            "special_direction"
                        ],
                        4,
                    ),
                )
                plt.plot(
                    magnitudes[order],
                    Energies_pred[order] - ref_omega_pace,
                    c=colors[i],
                )
                # plt.legend(bbox_to_anchor=(0.5, -0.05))
                plt.xlabel("Strain Magnitude")
                plt.ylabel(r"$\Delta E \enspace (eV/atom)$")
                i += 1
            plt.savefig("Strains_" + proto + "_pace.png", dpi=200)
            plt.close()


def plotlinescan(datapace: np.ndarray, datadft: np.ndarray, name: str, axisname: str):
    """
    plot stacking fault energy linescan
    input datapace
    """
    with plt.style.context("science"):
        plt.figure(figsize=(5, 4))
        plt.scatter(
            np.arange(0, 1.01, 1 / 12), datadft, marker=".", color="k", label="dft"
        )
        plt.plot(np.arange(0, 1.01, 1 / 12), datapace, color="k", label="pace")
        plt.xlabel(axisname)
        plt.ylabel(r"$Excess \enspace Energy \enspace (mJ/m^2)$")
        plt.legend()
        plt.tight_layout()
        plt.savefig(name)
        plt.close()


def get_whichlines() -> dict:
    whichlines = {}
    whichlines["bcc100"] = [("0_", r"$[100]$"), ("__", r"$[110]$")]
    whichlines["bcc110"] = [
        ("0_", r"$1/2 \enspace [1 \overline{1} 1]$"),
        ("__", r"$[100]$"),
        ("_-_", r"$[0 \overline{1} 1]$"),
    ]
    whichlines["fcc100"] = [("0_", r"$1/2 \enspace [110]$"), ("__", r"$[010]$")]
    whichlines["fcc111"] = [
        ("0_", r"$[101]$"),
        ("__", r"$[112]$"),
        ("_-_", r"$[1 \overline{1} 0]$"),
    ]
    whichlines["basal"] = [
        ("0_", r"$1/3 \enspace [1 1 \overline{2} 0]$"),
        ("__", r"$1/3 \enspace [\overline{2} 4 \overline{2} 0]$"),
        ("_-_", r"$1/3 \enspace [4 2 \overline{2} 0]$"),
    ]
    whichlines["prismatic"] = [
        ("0_", r"$1/3 \enspace [1 1 \overline{2} 0]$"),
        ("_0", r"$[0001]$"),
        ("__", r"$1/3 \enspace [1 1 \overline{2} 3]$"),
    ]
    whichlines["pyramidal"] = [
        ("_0", r"$1/3 \enspace [\overline{2} 1 1 0]$"),
        ("special", r"$1/3 \enspace [0 \overline{1} 1 2]$"),
    ]
    whichlines["bcc111"] = [
        ("0_", r"$[\overline{1} 1 0]$"),
        # ("__s", r"$[\overline{1} 0 1]$"),
        # ("_-_", r"$[\overline{1} 2 \overline{1}]$"),
    ]
    whichlines["pyramidal2nd"] = [
        ("0_", r"$[0 0 1 1]$"),
        # ("6_", r"$[0 0 1 1]$"),
    ]
    return whichlines


#####modified functions from uber.py#####
def shift_energies_to_surface_energy(unwinded, eqe):
    gamma = surface_energy(unwinded, eqe)
    unwinded["energy"] -= 2 * gamma + eqe
    return unwinded


def surface_energy(unwinded, eqe) -> float:
    maxcleave = max(unwinded["cleavage"])

    surface_energies = (
        unwinded.loc[unwinded["cleavage"] == maxcleave]["energy"] - eqe
    ) * 0.5
    if max(surface_energies) - min(surface_energies) > 0.001:
        raise ValueError(
            "Surface energies don't match at different shifts. Did you separate enough?"
        )

    return np.mean(surface_energies)


#            shift     frac    set_aspect_val                  xtext                                 ytext                                   text_pos
# bcc100       0      0.0385        1                        r'$[100]$'                            r'$[010]$'                    0.44,-0.1,-0.14,0.48,90,1.12,1.00
# bcc110      1/3     0.044    np.sqrt(2)*2/3   r'$1/2 \enspace [1 \overline{1} 1]$'   r'$1/2 \enspace [1 1 \overline{1}]$'      0.4,-0.12,-0.35,0.33,109.5,1.1,1.05
# fcc100       0      0.040         1           r'$1/2 \enspace [1 1 0]$'              r'$1/2 \enspace [\overline{1} 1 0]$'      0.4,-0.1,-0.14,0.4,90,1.12,1.00
# fcc111      0.5     0.041     np.sqrt(3)/2                 r'$[101]$'                            r'$[011]$'                    0.42,-0.13,-0.44,0.36,120,1.12,1.03
# basal       0.5     0.041     np.sqrt(3)/2    r'$1/3 \enspace [1 1 \overline{2} 1]$'       r'$[\overline{1} 1 0 0]$'           0.38,-0.13,-0.44,0.33,120,1.12,1.03
# prismatic    0      0.037        1.585        r'$1/3 \enspace [1 1 \overline{2} 0]$'            r'$[0001]$'                    0.32,-0.1,-0.18,0.45,90,1.18,1.00
# pyramidal  0.1427   0.024       0.5146        r'$1/3 \enspace [1 \overline{2} 1 3]$' r'$1/3 \enspace [\overline{2} 1 1 0]$'    0.38,-0.2,-0.24,0.23,105.5,1.08,1.02


def plot_interpolated(
    name: str,
    values: np.ndarray,
    shift: float,
    frac: float,
    set_aspect_val: float,
    xtext: str,
    ytext: str,
    pos: list,
    tpos: list,
    withline: tuple = None,
):
    """
    plot gsfe interpolated plots
    input:
        name: name of the face, for the name of the output picture
        values: gamma-surface Excess energies
        shift: shift the values according x-y angle
        frac: scaling for colorbar
        set_aspect_val: scaling for axis
        pos: x_coor_pos (2 nums),y_coor_pos (2 nums + 1 rotation),colorbar_pos (2 nums)
        tpos: tick positions
        withline: plot the linescan indication tuple with the name and the line position on picture
    """
    font30 = {
        "family": "Times New Roman",
        "weight": "normal",
        "size": 30,
    }
    font34 = {
        "family": "Times New Roman",
        "weight": "normal",
        "size": 34,
    }
    plt.rcParams["xtick.direction"] = "in"
    plt.rcParams["ytick.direction"] = "in"
    x = np.arange(0, 1.01, 1 / 12)
    y = np.arange(0, 1.01, 1 / 12)
    xx = np.arange(0, 1 + 1e-4, 1e-3)
    yy = np.arange(0, 1 + 1e-4, 1e-3)
    X, Y = np.meshgrid(xx, yy)
    Xnew = X - Y * shift
    Ynew = Y
    f = interpolate.interp2d(x, y, values, kind="cubic")
    values_new = f(xx, yy)
    fig = plt.figure(figsize=(15, 12))
    ax = fig.add_subplot(111)
    ax.spines["bottom"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["top"].set_visible(False)
    cont = ax.pcolormesh(Xnew, Ynew, values_new)
    im_ratio = (max(Ynew.ravel()) - min(Ynew.ravel())) / (
        max(Xnew.ravel()) - min(Xnew.ravel())
    )
    cbar = plt.colorbar(cont, ax=ax, fraction=frac * im_ratio, pad=0.04)
    cbar.ax.tick_params(labelsize=25, width=2, length=4)
    ax.set_xticks([])
    ax.set_yticks([])
    ax.set_ylim([-0.01, 1.01])
    # ax.set_xlim([-0.01,1.01])

    plt.plot([-shift, 0], [1, 0], c="k", linewidth=3)
    plt.plot([1 - shift, 1], [1, 0], c="k", linewidth=3)
    plt.plot([-shift, 1 - shift], [1, 1], c="k", linewidth=3)
    plt.plot([0, 1], [0, 0], c="k", linewidth=3)

    for i in [0.2, 0.4, 0.6, 0.8, 1.0]:
        length = 0.01
        start = (-i * shift, i)
        end = (-i * shift + length, i + length * shift)
        plt.plot([start[0], end[0]], [start[1], end[1]], linewidth=3, c="k")
        plt.text(start[0] - tpos[2], start[1] - tpos[3], str(i), font=font34)
        start = (i, 0)
        end = (i, length)
        plt.plot([start[0], end[0]], [start[1], end[1]], linewidth=3, c="k")
        plt.text(start[0] - tpos[0], start[1] - tpos[1], str(i), font=font34)
    start = (0, 0)
    end = (0, length)
    plt.plot([start[0], end[0]], [start[1], end[1]], linewidth=3, c="k")
    plt.text(start[0] - tpos[0], start[1] - tpos[1], str(0.0), font=font34)

    plt.text(pos[0], pos[1], xtext, font=font34, rotation=0)
    plt.text(pos[2], pos[3], ytext, font=font34, rotation=pos[4])
    plt.text(pos[5], pos[6], r"$mJ/m^2$", font=font30)

    for i in [0.1, 0.3, 0.5, 0.7, 0.9]:
        length = 0.005
        start = (-i * shift, i)
        end = (-i * shift + length, i + length * shift)
        plt.plot([start[0], end[0]], [start[1], end[1]], linewidth=2, c="k")
        start = (i, 0)
        end = (i, length)
        plt.plot([start[0], end[0]], [start[1], end[1]], linewidth=2, c="k")

    plt.gca().set_aspect(set_aspect_val)
    # plt.tight_layout()
    if withline == None:
        plt.savefig("interpolate_" + name + "_pace.png")
        plt.close()
    else:
        plt.plot(
            [
                withline[1][0] - withline[1][1] * shift,
                withline[1][2] - withline[1][3] * shift,
            ],
            [withline[1][1], withline[1][3]],
            linewidth=2,
            c="r",
        )
        if withline[1] == [0, 0, 1, 0.5]:
            plt.plot([0 - 0.5 * shift, 1 - 1 * shift], [0.5, 1], linewidth=2, c="r")
        if withline[1] == [0.5, 0, 1, 0.5]:
            plt.plot([0 - 0.5 * shift, 0.5-1 * shift], [0.5, 1], linewidth=2, c="r")
        plt.savefig("interp_" + name + withline[0] + ".png")

###### uber fit (only positive cleavages) #####
def uber_fit_positive(unwinded_slice):
    unwinded_slice = unwinded_slice.sort_values(["cleavage"])
    sigma=np.ones(len(unwinded_slice)-3)

    popt, pcov = curve_fit(uber, unwinded_slice["cleavage"][3:], unwinded_slice["energy"][3:],sigma=sigma,absolute_sigma=True)

    return popt,pcov

####### Generalized stacking fault energy #######
################ with uberfit ###################
def gsfe(data_collection: list, withUber: bool = True,fit_positive=False):
    """
    Generalized stacking fault energy analysis with UBER fittting
    input data_collection::list
    output pictures of the gamma-surfaces
    """
    faces = ["bcc100", "bcc110", "fcc100", "fcc111", "basal", "prismatic", "pyramidal","bcc111","pyramidal2nd"]
    protos = ["bcc", "bcc", "fcc", "fcc", "hcp", "hcp", "hcp","bcc","hcp"]
    planes = ["100", "110", "100", "111", "basal", "prismatic", "pyramidal","111","pyramidal2nd"]

    shifts = [0, 1 / 3, 0, 0.5, 0.5, 0, 0.1427,0.5,0]
    fracs = [0.0385, 0.044, 0.040, 0.041, 0.041, 0.039, 0.024,0.047,0.037]
    aspect_vals = [
        1,
        np.sqrt(2) * 2 / 3,
        1,
        np.sqrt(3) / 2,
        np.sqrt(3) / 2,
        1.585,
        0.5146,
        np.sqrt(3)/2,
        0.924,
    ]
    xtexts = [
        r"$[100]$",
        r"$1/2 \enspace [1 \overline{1} 1]$",
        r"$1/2 \enspace [1 1 0]$",
        r"$[101]$",
        r"$1/3 \enspace [1 1 \overline{2} 1]$",
        r"$1/3 \enspace [1 1 \overline{2} 0]$",
        r"$1/3 \enspace [1 \overline{2} 1 3]$",
        r'$[\overline{1} 1 0]$',
        r'$[0 0 1 1]$'
    ]
    ytexts = [
        r"$[010]$",
        r"$1/2 \enspace [1 1 \overline{1}]$",
        r"$1/2 \enspace [\overline{1} 1 0]$",
        r"$[011]$",
        r"$[\overline{1} 1 0 0]$",
        r"$[0001]$",
        r"$1/3 \enspace [\overline{2} 1 1 0]$",
        r'$[0 \overline{1} 1]$',
        r'$[ \overline{1} 1 0 0]$'
    ]
    # xlabel x y yalbel x y rot colorbar
    text_poses = [
        [0.44, -0.1, -0.14, 0.48, 90, 1.12, 1.00],
        [0.4, -0.12, -0.35, 0.33, 109.5, 1.1, 1.05],
        [0.4, -0.1, -0.14, 0.4, 90, 1.12, 1.00],
        [0.42, -0.13, -0.44, 0.36, 120, 1.12, 1.04],
        [0.38, -0.15, -0.46, 0.33, 120, 1.12, 1.03],
        [0.32, -0.12, -0.19, 0.45, 90, 1.18, 1.00],
        [0.38, -0.2, -0.24, 0.23, 105.5, 1.08, 1.04],
        [0.42, -0.13, -0.44,0.36, 120, 1.12,1.03],
        [0.32,-0.1, -0.18,0.45, 90, 1.18,1.00],
    ]
    # xtickshift ytickshift
    tick_poses = [
        [0.04, 0.05, 0.08, 0.02],
        [0.04, 0.05, 0.08, 0.03],
        [0.04, 0.05, 0.08, 0.03],
        [0.04, 0.06, 0.09, 0.03],
        [0.04, 0.06, 0.09, 0.04],
        [0.04, 0.05, 0.11, 0.02],
        [0.04, 0.09, 0.07, 0.03],
        [0.04, 0.05, 0.08, 0.03],
        [0.04, 0.05, 0.10, 0.02],

    ]

    whichlines = get_whichlines()
    for whichface in np.arange(9):
        if whichface!=8:
            continue
        face = faces[whichface]
        with open(face + "record.json") as f:
            shift_record = json.loads(f.read())
        equivalents = []
        for equi1 in shift_record["equivalents"]:
            equi_orbit = []
            for equi2 in equi1:
                if equi2.endswith("0.000000"):
                    a, b, c = equi2.split(":")
                    equi_orbit.append(a + "." + b)
                    # print(equi_orbit)
            if not equi_orbit == []:
                equivalents.append(equi_orbit)

        values_pace = np.zeros((13, 13), dtype=float)
        values_dft = np.zeros((13, 13), dtype=float)
        eqe_pace = 0
        eqe_dft = 0
        for data in data_collection:
            if (
                data["calc"] == "final"
                and data["metadata"]["perturbation"] == "gsfe"
                and data["metadata"]["proto"] == protos[whichface]
                and data["metadata"]["plane"] == planes[whichface]
                and data["metadata"]["shifts"][0] == 0
                and data["metadata"]["shifts"][1] == 0
                and data["metadata"]["cleavage"] == 0
                and not "relax" in data["metadata"]
            ):
                eqe_pace = data["pace"]["energy"]
                eqe_dft = data["energy"]
        for i in np.arange(len(equivalents)):
            first = equivalents[i][0]
            cleaves = []
            cleave_energies_pace = []
            cleave_energies_dft = []
            for data in data_collection:
                if (
                    data["calc"] == "final"
                    and data["metadata"]["perturbation"] == "gsfe"
                    and data["metadata"]["proto"] == protos[whichface]
                    and data["metadata"]["plane"] == planes[whichface]
                    and data["metadata"]["shifts"][0] == int(first.split(".")[0])
                    and data["metadata"]["shifts"][1] == int(first.split(".")[1])
                    and not "relax" in data["metadata"]
                ):
                    cleaves.append(data["metadata"]["cleavage"])
                    cleave_energies_pace.append(data["pace"]["energy"])
                    cleave_energies_dft.append(data["energy"])
            value_pace = 0
            value_dft = 0
            cleaves = np.array(cleaves)
            cleave_energies_pace = np.array(cleave_energies_pace)
            cleave_energies_dft = np.array(cleave_energies_dft)
            index = np.argsort(cleaves)
            energies_pace = cleave_energies_pace[index]
            energies_dft = cleave_energies_dft[index]
            cleave_record = load_record("./" +face+ "_cleave_record.json")
            unwinded_dft = tabulize_record(cleave_record)
            unwinded_pace = tabulize_record(cleave_record)
            cleaves_original=np.array([-0.15, -0.1, -0.05, 0.0, 0.05, 0.1, 0.15, 0.2, 0.3, 0.4, 0.5, 0.6, 0.7, 0.8, 0.9, 1.0, 1.25, 1.5, 1.75, 2.0, 2.5, 3.0, 3.5, 4.0, 6.0, 8.0])
            # cleave_min=min(cleaves)
            # argcleave=np.where(cleaves_original==cleave_min)[0]
            # if argcleave>0:
            #     to_drop=['0:0:%.6f'%cleaves_original[_] for _ in np.arange(argcleave)]
            #     unwinded_dft=unwinded_dft.drop(to_drop)
            #     unwinded_pace=unwinded_pace.drop(to_drop)
            if len(cleaves)<26:
                to_drop=['0:0:%.6f'%cleaves_orig for cleaves_orig in cleaves_original if not cleaves_orig in cleaves]
                unwinded_dft=unwinded_dft.drop(to_drop)
                unwinded_pace=unwinded_pace.drop(to_drop)

            unwinded_dft=unwinded_dft.sort_values(["cleavage"])
            unwinded_pace=unwinded_pace.sort_values(["cleavage"])
            
            avec, bvec = in_plane_vectors(shift_record)
            energies_pace = [
                transform_energy_to_surface_energy(e, avec, bvec) for e in energies_pace
            ]
            energies_dft = [
                transform_energy_to_surface_energy(e, avec, bvec) for e in energies_dft
            ]
            # if not len(unwinded_dft)==len(energies_dft)==len(energies_pace):
            #     print(faces[whichface],'_',first,'_',len(unwinded_dft),'_',energies_dft,'_',energies_pace)
            # try:
            #     unwinded_pace["energy"] = energies_pace
            # except:
                # print(faces[whichface],'_',first,'_',argcleave.sort(),'_',len(unwinded_dft),'_',len(energies_dft),'_',len(energies_pace))
                # print(cleave_min)
                # print(cleaves)
            unwinded_pace["energy"] = energies_pace
            unwinded_dft["energy"] = energies_dft
            unwinded_slice_pace = shift_energies_to_surface_energy(
                unwinded_pace, eqe_pace
            )
            unwinded_slice_dft = shift_energies_to_surface_energy(unwinded_dft, eqe_dft)
            try:
                popt_pace, pcov_pace = uber_fit(unwinded_slice_pace)
                popt_dft, pcov_dft = uber_fit(unwinded_slice_dft)
            except:
                pass
            if fit_positive==True:
                popt_pace, pcov_pace = uber_fit_positive(unwinded_slice_pace)
                popt_dft, pcov_dft = uber_fit_positive(unwinded_slice_dft)                
            if withUber == True:
                value_pace = -2 * popt_pace[1] + popt_pace[3]
                value_dft = -2 * popt_dft[1] + popt_dft[3]
            else:
                value_pace = energies_pace[3]
                value_dft = energies_dft[3]
            # if (
            #     planes[whichface] == 'pyramidal2nd'
            #     and int(first.split(".")[0]) == 0
            #     and int(first.split(".")[1]) == 5
            # ):
            #     print(value_pace)
            #     print(value_dft)
            
            if planes[whichface] == 'pyramidal2nd' and int(first.split(".")[0]) == 0 and int(first.split(".")[1]) == 5:
                print(avec)
                print(bvec)
                fig = plt.figure()
                ax = fig.add_subplot(111)
                unwinded_slice_pace=unwinded_slice_pace.sort_values(["cleavage"])
                print(unwinded_slice_pace["cleavage"])
                deltas = np.linspace(
                    min(unwinded_slice_pace["cleavage"]),
                    max(unwinded_slice_pace["cleavage"]),
                    1000,
                )
                fitenergy = uber(deltas, *popt_pace)
                ax.plot(deltas, fitenergy - popt_pace[3], "g--", color="k", zorder=1)
                ax.scatter(
                    unwinded_slice_pace["cleavage"],
                    unwinded_slice_pace["energy"] - popt_pace[3],
                    zorder=2,
                )
                ax.set_xlabel("$d$ $\mathrm{[\AA]}$")
                ax.set_ylabel("Energy $\mathrm{[J/m^2]}$")
                plt.tight_layout()
                plt.savefig(
                    "Uberplot"
                    + first.split(".")[0]
                    + "_"
                    + first.split(".")[1]
                    + "_pace.png"
                )
                plt.close()
                fig = plt.figure()
                ax = fig.add_subplot(111)
                unwinded_slice_dft=unwinded_slice_dft.sort_values(["cleavage"])
                print(unwinded_slice_dft["cleavage"])
                deltas = np.linspace(
                    min(unwinded_slice_dft["cleavage"]),
                    max(unwinded_slice_dft["cleavage"]),
                    1000,
                )
                fitenergy = uber(deltas, *popt_dft)
                ax.plot(deltas, fitenergy - popt_dft[3], "g--", color="k", zorder=1)
                ax.scatter(
                    unwinded_slice_dft["cleavage"],
                    unwinded_slice_dft["energy"] - popt_dft[3],
                    zorder=2,
                )
                ax.set_xlabel("$d$ $\mathrm{[\AA]}$")
                ax.set_ylabel("Energy $\mathrm{[J/m^2]}$")
                plt.tight_layout()
                plt.savefig(
                    "Uberplot"
                    + first.split(".")[0]
                    + "_"
                    + first.split(".")[1]
                    + "_dft.png"
                )
                plt.close()

            for j in np.arange(len(equivalents[i])):
                which = equivalents[i][j].split(".")
                values_pace[int(which[0]), int(which[1])] = value_pace
                values_dft[int(which[0]), int(which[1])] = value_dft

        values_pace[12, :] = values_pace[0, :]
        values_pace[:, 12] = values_pace[:, 0]
        values_pace = (values_pace - values_pace[0, 0]) * 1000

        values_dft[12, :] = values_dft[0, :]
        values_dft[:, 12] = values_dft[:, 0]
        values_dft = (values_dft - values_dft[0, 0]) * 1000

        def plot_surface_heatmap(ax, X, Y, Z):
            pc = ax.pcolormesh(X, Y, Z)
            ax.set_aspect("equal")
            im_ratio = (max(Y.ravel()) - min(Y.ravel())) / (
                max(X.ravel()) - min(X.ravel())
            )
            cbar = plt.colorbar(pc, ax=ax, fraction=0.047 * im_ratio, pad=0.04)

            return ax, cbar

        fig = plt.figure()
        ax = fig.add_subplot(111)
        x = np.arange(0, 1 + 1e-7, 1 / 12)
        y = np.arange(0, 1 + 1e-7, 1 / 12)
        X, Y = np.meshgrid(x, y)
        ax, cbar = plot_surface_heatmap(ax, X, Y, values_pace)
        plt.tight_layout()
        plt.savefig("2d_" + faces[whichface] + "_pace.png")
        plt.close()

        plot_interpolated(
            faces[whichface],
            values_pace,
            shifts[whichface],
            fracs[whichface],
            aspect_vals[whichface],
            xtexts[whichface],
            ytexts[whichface],
            text_poses[whichface],
            tick_poses[whichface],
        )

        for line in whichlines[faces[whichface]]:
            if line[0] == "0_":
                plotlinescan(
                    values_pace[0, :],
                    values_dft[0, :],
                    faces[whichface] + "-" + "linescan0_.png",
                    line[1],
                )
                plot_interpolated(
                    faces[whichface],
                    values_pace,
                    shifts[whichface],
                    fracs[whichface],
                    aspect_vals[whichface],
                    xtexts[whichface],
                    ytexts[whichface],
                    text_poses[whichface],
                    tick_poses[whichface],
                    (line[0], [0, 0.05, 1, 0.05]),
                )
            elif line[0] == "6_":
                plotlinescan(
                    values_pace[6, :],
                    values_dft[6, :],
                    faces[whichface] + "-" + "linescan6_.png",
                    line[1],
                )
                plot_interpolated(
                    faces[whichface],
                    values_pace,
                    shifts[whichface],
                    fracs[whichface],
                    aspect_vals[whichface],
                    xtexts[whichface],
                    ytexts[whichface],
                    text_poses[whichface],
                    tick_poses[whichface],
                    (line[0], [0, 0.5, 1, 0.5]),
                )
            elif line[0] == "_0":
                plotlinescan(
                    values_pace[:, 0],
                    values_dft[:, 0],
                    faces[whichface] + "-" + "linescan_0.png",
                    line[1],
                )
                plot_interpolated(
                    faces[whichface],
                    values_pace,
                    shifts[whichface],
                    fracs[whichface],
                    aspect_vals[whichface],
                    xtexts[whichface],
                    ytexts[whichface],
                    text_poses[whichface],
                    tick_poses[whichface],
                    (line[0], [0.05, 0, 0.05, 1]),
                )
            elif line[0] == "__":
                plotlinescan(
                    np.diag(values_pace),
                    np.diag(values_dft),
                    faces[whichface] + "-" + "linescan__.png",
                    line[1],
                )
                plot_interpolated(
                    faces[whichface],
                    values_pace,
                    shifts[whichface],
                    fracs[whichface],
                    aspect_vals[whichface],
                    xtexts[whichface],
                    ytexts[whichface],
                    text_poses[whichface],
                    tick_poses[whichface],
                    (line[0], [0, 0, 1, 1]),
                )
            elif line[0] == "__s":
                plotlinescan(
                    values_pace[
                        np.array([6,7,8,9,10,11,12,0,1,2,3,4,5]),
                        np.arange(13),
                    ],
                    values_dft[
                        np.array([6,7,8,9,10,11,12,0,1,2,3,4,5]),
                        np.arange(13),
                    ],
                    faces[whichface] + "-" + "linescan__s.png",
                    line[1],
                )
                plot_interpolated(
                    faces[whichface],
                    values_pace,
                    shifts[whichface],
                    fracs[whichface],
                    aspect_vals[whichface],
                    xtexts[whichface],
                    ytexts[whichface],
                    text_poses[whichface],
                    tick_poses[whichface],
                    (line[0], [0.5,0, 1, 0.5]),
                )
            elif line[0] == "_-_":
                plotlinescan(
                    np.diag(np.flipud(values_pace)),
                    np.diag(np.flipud(values_dft)),
                    faces[whichface] + "-" + "linescan_-_.png",
                    line[1],
                )
                plot_interpolated(
                    faces[whichface],
                    values_pace,
                    shifts[whichface],
                    fracs[whichface],
                    aspect_vals[whichface],
                    xtexts[whichface],
                    ytexts[whichface],
                    text_poses[whichface],
                    tick_poses[whichface],
                    (line[0], [0, 1, 1, 0]),
                )
            elif line[0] == "special":
                plotlinescan(
                    values_pace[
                        np.arange(13),
                        np.array([0, 2, 4, 6, 8, 10, 12, 2, 4, 6, 8, 10, 12]),
                    ],
                    values_dft[
                        np.arange(13),
                        np.array([0, 2, 4, 6, 8, 10, 12, 2, 4, 6, 8, 10, 12]),
                    ],
                    faces[whichface] + "-" + "linescanspecial.png",
                    line[1],
                )
                plot_interpolated(
                    faces[whichface],
                    values_pace,
                    shifts[whichface],
                    fracs[whichface],
                    aspect_vals[whichface],
                    xtexts[whichface],
                    ytexts[whichface],
                    text_poses[whichface],
                    tick_poses[whichface],
                    (line[0], [0, 0, 1, 0.5]),
                )


def burgers_bains(data_collection: list):
    """
    analysis for Burger_Bains pathways connecting fcc-bcc-hcp
    """
    B_B_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "pathways"
            and data_collection[i]["calc"] == "final"
        ):
            if data_collection[i]['metadata']['pathway_label']=="Burgers_Bain":
                B_B_id.append(i)
    bccref_pace = 0
    bccref_dft = 0
    for data in data_collection:
        if (
            data["metadata"]["perturbation"] == "icsd"
            and data["calc"] == "final"
            and data["metadata"]["proto"] == "bcc"
        ):
            bccref_pace = data["pace"]["energy"] / len(data["structure"])
            bccref_dft = data["energy"] / len(data["structure"])

    bccref_pace
    unshuffled_values = []
    unshuffled_energies_dft = []
    unshuffled_energies_pace = []

    shuffled_values = []
    shuffled_energies_dft = []
    shuffled_energies_pace = []

    for id in B_B_id:
        if data_collection[id]["metadata"]["shuffled"]==True:
            shuffled_values.append(
                data_collection[id]["metadata"]["value"]
            )
            shuffled_energies_dft.append(
                float(data_collection[id]["energy"])
                / len(data_collection[id]["structure"])
                - bccref_dft
            )
            shuffled_energies_pace.append(
                float(data_collection[id]["pace"]["energy"])
                / len(data_collection[id]["structure"])
                - bccref_pace
            )
        else:
            unshuffled_values.append(
                data_collection[id]["metadata"]["value"]
            )
            unshuffled_energies_dft.append(
                float(data_collection[id]["energy"])
                / len(data_collection[id]["structure"])
                - bccref_dft
            )
            unshuffled_energies_pace.append(
                float(data_collection[id]["pace"]["energy"])
                / len(data_collection[id]["structure"])
                - bccref_pace
            )
    unshuffled_values = np.array(unshuffled_values)
    unshuffled_energies_dft = np.array(unshuffled_energies_dft)
    unshuffled_energies_pace = np.array(unshuffled_energies_pace)

    shuffled_values = np.array(shuffled_values)
    shuffled_energies_dft = np.array(shuffled_energies_dft)
    shuffled_energies_pace = np.array(shuffled_energies_pace)

    unshuffled_index = np.argsort(unshuffled_values)
    shuffled_index = np.argsort(shuffled_values)

    unshuffled_values = unshuffled_values[unshuffled_index]
    unshuffled_energies_dft = unshuffled_energies_dft[unshuffled_index]
    unshuffled_energies_pace = unshuffled_energies_pace[unshuffled_index]

    shuffled_values = shuffled_values[shuffled_index]
    shuffled_energies_dft = shuffled_energies_dft[shuffled_index]
    shuffled_energies_pace = shuffled_energies_pace[shuffled_index]

    # print(shuffled_values)
    with plt.style.context("science"):
        plt.figure(figsize=(5, 4))
        plt.scatter(
            shuffled_values, shuffled_energies_dft, s=6, c="r", label="shuffled_dft"
        )
        plt.scatter(
            unshuffled_values,
            unshuffled_energies_dft,
            s=6,
            c="b",
            label="unshuffled_dft",
        )
        plt.plot(shuffled_values, shuffled_energies_pace, c="r", label="shuffled_pace")
        plt.plot(
            unshuffled_values, unshuffled_energies_pace, c="b", label="unshuffled_pace"
        )
        plt.xlabel("Strain Magnitude")
        plt.ylabel(r"$\Delta E \enspace (eV/atom)$")
        plt.tick_params(axis="y", direction="in")
        plt.tick_params(axis="x", direction="in")
        plt.legend()
        plt.savefig("Ti_Burgers_strain.png")
        plt.close()


def bcc_omega(data_collection: list, ref_omega_dft: float, ref_omega_pace: float):
    """
    analysis for the pathway connecting bcc and omega
    ref_omega_dft::float   reference energy (per atom) for omega structure from dft calculation
    ref_omega_pace::float  reference energy (per atom) for omega structure from pace potential
    """
    b_o_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "pathways"
            and data_collection[i]["calc"] == "final"
        ):
            if data_collection[i]["metadata"]["pathway_label"] == "bcc_omega":
                b_o_id.append(i)
    values = []
    Energies_dft = []
    Energies_pace = []
    for id in b_o_id:
        values.append(data_collection[id]["metadata"]["value"])
        Energies_pace.append(
            data_collection[id]["pace"]["energy"]
            / len(data_collection[id]["structure"])
        )
        Energies_dft.append(
            data_collection[id]["energy"] / len(data_collection[id]["structure"])
        )
    values = np.array(values)
    Energies_dft = np.array(Energies_dft)
    Energies_pace = np.array(Energies_pace)
    index = np.argsort(values)
    values = values[index]
    Energies_dft = Energies_dft[index]
    Energies_pace = Energies_pace[index]
    with plt.style.context("science"):
        plt.figure()
        plt.scatter(
            np.arange(len(Energies_dft)),
            Energies_dft - ref_omega_dft,
            c="black",
            s=4,
            label="dft",
        )
        plt.plot(
            np.arange(len(Energies_pace)),
            Energies_pace - ref_omega_pace,
            c="black",
            label="pace",
        )
        plt.xlabel("Omega Transformation")
        plt.xticks([])
        plt.ylabel(r"$\Delta E^{pred} \enspace (eV/atom)$")
        plt.legend()
        plt.savefig("Omega.png", dpi=150)
        plt.close()


def hcp_omega(data_collection: list):
    """
    analysis for the pathways connecting hcp and omega
    """
    h_o_static_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "pathways"
            and data_collection[i]["metadata"]["pathway_label"] == "hcp_omega"
            and data_collection[i]["metadata"]["SSNEB"] == False
        ):
            h_o_static_id.append(i)

    h_o_NEB_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "pathways"
            and data_collection[i]["metadata"]["pathway_label"] == "hcp_omega"
            and data_collection[i]["metadata"]["SSNEB"] == True
        ):
            h_o_NEB_id.append(i)
    hcp_energy_pace = 0
    omega_energy_pace = 0
    hcp_energy_dft = 0
    omega_energy_dft = 0
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "icsd"
            and data_collection[i]["metadata"]["proto"] == "hcp"
            and data_collection[i]["calc"] == "final"
        ):
            hcp_energy_dft = data_collection[i]["energy"] / len(
                data_collection[i]["structure"]
            )
            hcp_energy_pace = data_collection[i]["pace"]["energy"] / len(
                data_collection[i]["structure"]
            )
        if (
            data_collection[i]["metadata"]["perturbation"] == "icsd"
            and data_collection[i]["metadata"]["proto"] == "omega"
            and data_collection[i]["calc"] == "final"
        ):
            omega_energy_dft = data_collection[i]["energy"] / len(
                data_collection[i]["structure"]
            )
            omega_energy_pace = data_collection[i]["pace"]["energy"] / len(
                data_collection[i]["structure"]
            )
    pathways = set()
    for id in h_o_static_id:
        pathways.add(data_collection[id]["metadata"]["pathway_number"])
    pathways = list(pathways)

    for pathway in pathways:
        energies_dft = np.zeros(7)
        energies_dft[0] = omega_energy_dft
        energies_dft[6] = hcp_energy_dft
        energies_pace = np.zeros(7)
        energies_pace[0] = omega_energy_pace
        energies_pace[6] = hcp_energy_pace
        for image in np.arange(1, 6, 1):
            for i in np.arange(len(data_collection)):
                if (
                    "SSNEB" in data_collection[i]["metadata"]
                    and data_collection[i]["metadata"]["perturbation"] == "pathways"
                    and data_collection[i]["metadata"]["pathway_label"] == "hcp_omega"
                    and data_collection[i]["metadata"]["SSNEB"] == False
                    and data_collection[i]["metadata"]["pathway_number"] == pathway
                    and int(data_collection[i]["metadata"]["image"]) == image
                ):
                    energies_pace[image] = data_collection[i]["pace"]["energy"] / len(
                        data_collection[i]["structure"]
                    )
                    energies_dft[image] = data_collection[i]["energy"] / len(
                        data_collection[i]["structure"]
                    )
        energies_dft = energies_dft - energies_dft[0]
        energies_pace = energies_pace - energies_pace[0]

        energies_dft_NEB = np.zeros(7)
        energies_dft_NEB[0] = omega_energy_dft
        energies_dft_NEB[6] = hcp_energy_dft
        energies_pace_NEB = np.zeros(7)
        energies_pace_NEB[0] = omega_energy_pace
        energies_pace_NEB[6] = hcp_energy_pace
        for image in np.arange(1, 6, 1):
            for i in np.arange(len(data_collection)):
                if (
                    "SSNEB" in data_collection[i]["metadata"]
                    and data_collection[i]["metadata"]["perturbation"] == "pathways"
                    and data_collection[i]["metadata"]["pathway_label"] == "hcp_omega"
                    and data_collection[i]["metadata"]["SSNEB"] == True
                    and data_collection[i]["metadata"]["pathway_number"] == pathway
                    and int(data_collection[i]["metadata"]["image"]) == image
                ):
                    energies_pace_NEB[image] = data_collection[i]["pace"][
                        "energy"
                    ] / len(data_collection[i]["structure"])
                    energies_dft_NEB[image] = data_collection[i]["energy"] / len(
                        data_collection[i]["structure"]
                    )
        energies_dft_NEB = energies_dft_NEB - energies_dft_NEB[0]
        energies_pace_NEB = energies_pace_NEB - energies_pace_NEB[0]
        with plt.style.context("science"):
            plt.figure(figsize=(5, 4))
            plt.scatter(np.arange(7), energies_dft, c="k", label="dft")
            plt.plot(np.arange(7), energies_pace, c="k", label="pace")
            plt.scatter(np.arange(7), energies_dft_NEB, c="r", label="dft_SSNEB")
            plt.plot(np.arange(7), energies_pace_NEB, c="r", label="pace_SSNEB")
            plt.xlabel("Transformation Coordinates")
            plt.ylabel(r"$\Delta E \enspace (eV/atom)$")
            plt.legend()
            plt.xticks([0, 6], ["omega", "hcp"])
            plt.savefig(str(pathway) + "_pace.png")
            plt.close()


def calculate_phonons(calc_list: list,pace=False) -> Phonopy:
    phonopy_supercell = get_phonopy_structure(calc_list[0]["metadata"]["reference_cell"])
    phonon = Phonopy(
        unitcell=phonopy_supercell,
        primitive_matrix="auto",
    )
    # phonon.generate_displacements(distance=0.03)
    # find the order in which the displacements are enumerated within the phonopy
    # object

    phonon_dataset = {"natom": len(calc_list[0]["metadata"]["reference_cell"]), "first_atoms": []}
    for calc in calc_list:
        calc_dataset = calc["metadata"]["phonopy"]
        if pace==False:
            calc_dataset["forces"] = np.array(calc["forces"])
        else:
            calc_dataset["forces"] = np.array(calc["pace"]["forces"])
        phonon_dataset["first_atoms"].append(calc_dataset)
    phonon.dataset = phonon_dataset
    phonon.produce_force_constants()
    return phonon


def phonos(data_collection: list):
    """
    analysis for the phonon dispersion curves
    """
    phonon_proto=set()
    phonon_collection=[]
    for data in data_collection:
        if data["calc"] == "final" and data["metadata"]["perturbation"] == "phonons":
            phonon_proto.add(data["metadata"]['proto'])
            phonon_collection.append(data)
    phonon_calc_dict = dict()
    for data in phonon_collection:
        if data['metadata']['proto'] not in phonon_calc_dict:
            phonon_calc_dict[data['metadata']["proto"]] = []
        phonon_calc_dict[data['metadata']["proto"]].append(data)
    phonopy_dict = dict()
    for proto in phonon_calc_dict:
        phonopy_dict[proto] = calculate_phonons(phonon_calc_dict[proto])
        band_fig = phonopy_dict[proto].auto_band_structure(plot=True)
        for ax in band_fig.gcf().axes:
            ax.tick_params(axis="both", labelsize=12.0)
            for line_idx in range(len(ax.get_lines())):
                ax.get_lines()[line_idx].set_color("black")
        band_fig.gcf().axes[0].set_ylabel(r"Frequency (THz)", fontsize=14.0)
        band_fig.savefig(os.path.join("phonon_{0}.png".format(proto)), transparent=True)
    for proto in phonon_calc_dict:
        phonopy_dict[proto] = calculate_phonons(phonon_calc_dict[proto],pace=True)
        band_fig = phonopy_dict[proto].auto_band_structure(plot=True)
        for ax in band_fig.gcf().axes:
            ax.tick_params(axis="both", labelsize=12.0)
            for line_idx in range(len(ax.get_lines())):
                ax.get_lines()[line_idx].set_color("black")
        band_fig.gcf().axes[0].set_ylabel(r"Frequency (THz)", fontsize=14.0)
        band_fig.savefig(os.path.join("phonon_{0}_pace.png".format(proto)), transparent=True)






def collect_in_powerpoint(
    data_collection:list, RMSE_E: float, MAE_E: float, RMSE_F: float, MAE_F: float, pptxname:str="pace" 
):
    """
    collect things in a powerpoint file
    input   RMSE_E::float Energy RMSE error
            MAE_E::float   Energy MAE error
            RMSE_F::float  Force RMSE error
            MAE_F::float   Force MAE error
    """
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = ""
    p = tf.add_paragraph()
    p.text = "Energy"
    p.font.size = Pt(30)
    p0 = tf.add_paragraph()
    p0.text = ""
    p1 = tf.add_paragraph()
    p1.text = "RMSE=%.5f" % RMSE_E
    p2 = tf.add_paragraph()
    p2.text = "MAE=%.5f" % MAE_E
    img_path = "pace_energy.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    img_path = "pace_energy_low.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = ""
    p = tf.add_paragraph()
    p.text = "Force"
    p.font.size = Pt(30)
    p0 = tf.add_paragraph()
    p0.text = ""
    p1 = tf.add_paragraph()
    p1.text = "RMSE=%.5f" % RMSE_F
    p2 = tf.add_paragraph()
    p2.text = "MAE=%.5f" % MAE_F
    img_path = "pace_force.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    img_path = "pace_force_low.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    for i in np.arange(4):
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Volumetric_" + str(i)
        img_path = "Volumetric_deformation_" + str(i) + "_pace.png"
        picleft = Inches(3)
        picheight = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Volumetric deformation"
    img_path = "Volumetric_deformation_pace.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Grain Boundary"
    img_path = "Grain_boundary_pace.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Surfaces"
    img_path = "Surfaces_pace.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Vacancy formation"
    img_path = "Vacancy_formation_pace.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Burgers_Bains"
    img_path = "Ti_Burgers_strain.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Omega transaformation"
    img_path = "Omega.png"
    picleft = Inches(3)
    picheight = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    strain_id = []
    for i in np.arange(len(data_collection)):
        if (
            data_collection[i]["metadata"]["perturbation"] == "strain"
            and data_collection[i]["calc"] == "final"
        ):
            strain_id.append(i)

    protos = set()
    for id in strain_id:
        protos.add(data_collection[id]["metadata"]["proto"])
    protos = list(protos)
    for proto in protos:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        name=proto
        # if name in renaming:
        #     name = renaming[name]
        tf.text=name
        img_path = "Strains_" + proto + "_pace.png"
        picleft = Inches(3)
        picheight = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    pathnum = [
        "0",
        "1",
        "67",
        "297",
        "357",
        "499",
    ]

    for ptn in pathnum:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = ptn
        img_path = ptn + "_pace.png"
        picleft = Inches(3)
        picheight = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    faces = ["bcc100", "bcc110", "fcc100", "fcc111", "basal", "prismatic", "pyramidal","bcc111","pyramidal2nd"]
    for face in faces:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = face
        img_path = "2d_" + face + "_pace.png"
        picleft = Inches(3)
        picheight = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, picleft, top, height=picheight)

    linescans_ = glob.glob("*linescan*")
    #    {'bcc110-linescan_-_.png':,
    #     'bcc110-linescan_-_.png':,
    #     'prismatic-linescan__.png':,
    #     'fcc111-linescan0_.png';,
    #     'bcc110-linescan0_.png':,
    #     'fcc111-linescan_-_.png':,
    #     'basal-linescan0_.png':,
    #     'bcc100-linescan0_.png':,
    #     'fcc100-linescan0_.png':,
    #     'bcc100-linescan__.png':,
    #     'basal-linescan__.png':,
    #     'fcc100-linescan__.png':,
    #     'prismatic-linescan_0.png':,
    #     'pyramidal-linescan_-_.png':,
    #     'prismatic-linescan0_.png':,
    #     'bcc110-linescan__.png':,
    #     'basal-linescan_-_.png':,
    #     'pyramidal-linescan_0.png':,
    #     'fcc111-linescan__.png':}
    linescans1 = [l for l in linescans_ if "bcc100" in l]
    linescans2 = [l for l in linescans_ if "bcc110" in l]
    linescans3 = [l for l in linescans_ if "fcc100" in l]
    linescans4 = [l for l in linescans_ if "fcc111" in l]
    linescans5 = [l for l in linescans_ if "basal" in l]
    linescans6 = [l for l in linescans_ if "prismatic" in l]
    linescans7 = [l for l in linescans_ if "pyramidal" in l and "2nd" not in l]
    linescans8 = [l for l in linescans_ if "bcc111" in l]
    linescans9 = [l for l in linescans_ if "pyramidal2nd" in l]

    # linescans = (
    #     linescans1
    #     + linescans2
    #     + linescans3
    #     + linescans4
    #     + linescans5
    #     + linescans6
    #     + linescans7
    # )
    print(linescans7)
    facenames = [
        "bcc100",
        "bcc110",
        "fcc100",
        "fcc111",
        "basal",
        "prismatic",
        "pyramidal",
        "bcc111",
        "pyramidal2nd"
    ]
    names = locals()
    for i in np.arange(1, 10, 1):
        for linescan in names["linescans" + str(i)]:
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = linescan.split(".")[0]
            # img_path_face="interpolate_" + facenames[i-1] + "_pace.png"
            img_path_face = (
                "interp_"
                + facenames[i - 1]
                + linescan.split("scan")[-1].split(".png")[0]
                + ".png"
            )
            pic = slide.shapes.add_picture(
                img_path_face, Inches(0), Inches(2), height=Inches(3.5)
            )
            img_path = linescan
            pic = slide.shapes.add_picture(
                img_path, Inches(5), Inches(2), height=Inches(3.5)
            )
    phonon_proto=set()
    phonon_collection=[]
    for data in data_collection:
        if data["calc"] == "final" and data["metadata"]["perturbation"] == "phonons":
            phonon_proto.add(data["metadata"]['proto'])
    phonon_proto=list(phonon_proto)

    for proto in phonon_proto:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        name=proto
        # if name in renaming:
        #     name = renaming[name]
        tf.text = name
        if proto==phonon_proto[0]:
            p = tf.add_paragraph()
            p.text = "left:dft"
            p0 = tf.add_paragraph()
            p0.text = "right:pace"
        img_path = "phonon_"+proto+".png"
        pic = slide.shapes.add_picture(img_path, Inches(0), Inches(2), height=Inches(3.5))
        img_path = "phonon_"+proto+"_pace.png"
        pic = slide.shapes.add_picture(img_path, Inches(5), Inches(2), height=Inches(3.5))

    prs.save(pptxname+".pptx")


def main(pptxname:str='pace'):
    # filename=glob(r'*.json')[0]
    filename = "Ti_data_pace.json"
    with open(filename) as f:
        data_collection_ = json.loads(f.read(), cls=MontyDecoder)
    data_collection = []
    for data in data_collection_:
        if data["metadata"]["perturbation"] != "pairs":
            data_collection.append(data)
    del data_collection_

    ######get reference energies#####
    ref_omega_dft = 0
    for data in data_collection:
        if (
            data["metadata"]["perturbation"] == "icsd"
            and data["calc"] == "final"
            and data["metadata"]["proto"] == "omega"
        ):
            ref_omega_dft = data["energy"]
    ref_omega_dft = ref_omega_dft / 3

    ref_omega_pace = 0
    for data in data_collection:
        if (
            data["metadata"]["perturbation"] == "icsd"
            and data["calc"] == "final"
            and data["metadata"]["proto"] == "omega"
        ):
            ref_omega_pace = data["pace"]["energy"]
    ref_omega_pace = ref_omega_pace / 3

    #####evaluate energy and force prediction#####
    energies_dft = []
    forces_dft = []
    energies_pace = []
    forces_pace = []

    energies_dft_stdt = []
    forces_dft_stdt = []
    energies_pace_stdt = []
    forces_pace_stdt = []
    for i in tqdm(np.arange(len(data_collection))):
        if not "standout" in data_collection[i].keys():
            atom_num = len(data_collection[i]["structure"].species)
            energies_dft.append(data_collection[i]["energy"] / atom_num - ref_omega_dft)
            forces_dft += data_collection[i]["forces"]
            energies_pace.append(
                data_collection[i]["pace"]["energy"] / atom_num - ref_omega_pace
            )
            forces_pace += data_collection[i]["pace"]["forces"].tolist()
            if data_collection[i]["pace"]["energy"] / atom_num - ref_omega_pace>20:
                print(data_collection[i]['metadata'])
        else:
            atom_num = len(data_collection[i]["structure"].species)
            energies_dft_stdt.append(
                data_collection[i]["energy"] / atom_num - ref_omega_dft
            )
            forces_dft_stdt += data_collection[i]["forces"]
            energies_pace_stdt.append(
                data_collection[i]["pace"]["energy"] / atom_num - ref_omega_pace
            )
            forces_pace_stdt += data_collection[i]["pace"]["forces"].tolist()

    with plt.style.context("science"):
        plt.figure()
        plt.scatter(energies_dft, energies_pace, c="k", s=0.3, label="non-holdout")
        plt.scatter(
            energies_dft_stdt, energies_pace_stdt, c="r", s=0.3, label="holdout"
        )
        plt.xlabel(r"$\Delta E^{DFT} \enspace (eV/atom)$")
        plt.ylabel(r"$\Delta E^{pred} \enspace (eV/atom)$")
        # plt.legend()
        plt.savefig("pace_energy.png", dpi=200)
        plt.close()

    with plt.style.context("science"):
        plt.figure()
        plt.scatter(energies_dft, energies_pace, c="k", s=0.3, label="non-holdout")
        plt.scatter(
            energies_dft_stdt, energies_pace_stdt, c="r", s=0.3, label="holdout"
        )
        plt.xlabel(r"$\Delta E^{DFT} \enspace (eV/atom)$")
        plt.ylabel(r"$\Delta E^{pred} \enspace (eV/atom)$")
        plt.xlim([0, 2])
        plt.ylim([0, 2])
        # plt.legend()
        plt.savefig("pace_energy_low.png", dpi=200)
        plt.close()

    energies_dft += energies_dft_stdt
    energies_pace += energies_pace_stdt
    RMSE_E = np.sqrt(mean_squared_error(energies_dft, energies_pace))
    MAE_E = mean_absolute_error(energies_dft, energies_pace)

    forces_sca_dft = np.sqrt(np.sum(np.array(forces_dft) ** 2, axis=1))
    forces_sca_pace = np.sqrt(np.sum(np.array(forces_pace) ** 2, axis=1))
    try:
        forces_sca_dft_stdt = np.sqrt(np.sum(np.array(forces_dft_stdt) ** 2, axis=1))
        forces_sca_pace_stdt = np.sqrt(np.sum(np.array(forces_pace_stdt) ** 2, axis=1))
    except:
        pass

    with plt.style.context("science"):
        plt.figure()
        plt.scatter(forces_sca_dft, forces_sca_pace, c="k", s=0.3, label="non-holdout")
        try:
            plt.scatter(
                forces_sca_dft_stdt, forces_sca_pace_stdt, c="r", s=0.3, label="holdout"
            )
        except:
            pass
        plt.xlabel(r"$|F^{i}|^{DFT} \enspace (eV/Å)$")
        plt.ylabel(r"$|F^{i}|^{pred} \enspace (eV/Å)$")
        # plt.legend()
        plt.savefig("pace_force.png", dpi=200)
        plt.close()

    with plt.style.context("science"):
        plt.figure()
        plt.scatter(forces_sca_dft, forces_sca_pace, c="k", s=0.3, label="non-holdout")
        try:
            plt.scatter(
                forces_sca_dft_stdt, forces_sca_pace_stdt, c="r", s=0.3, label="holdout"
            )
        except:
            pass
        plt.xlabel(r"$|F^{i}|^{DFT} \enspace (eV/Å)$")
        plt.ylabel(r"$|F^{i}|^{pred} \enspace (eV/Å)$")
        # plt.legend()
        plt.xlim([0, 10])
        plt.ylim([0, 10])
        plt.savefig("pace_force_low.png", dpi=200)
        plt.close()

    try:
        forces_sca_dft += forces_sca_dft_stdt
        forces_sca_pace += forces_sca_pace_stdt
    except:
        pass

    RMSE_F = np.sqrt(mean_squared_error(forces_sca_dft, forces_sca_pace))
    MAE_F = mean_absolute_error(forces_sca_dft, forces_sca_pace)
    vacancy_formation(data_collection)
    print('Vacancy')
    volumetric_deformation(data_collection, ref_omega_dft, ref_omega_pace)
    print('Volumetric')
    grain_boundary(data_collection)
    print('Grain Boundary')
    surface(data_collection)
    print('Surface')
    gsfe(data_collection,withUber=False,fit_positive=False)
    print('Stacking Fault Energy')
    burgers_bains(data_collection)
    bcc_omega(data_collection, ref_omega_dft, ref_omega_pace)
    hcp_omega(data_collection)
    print('Pathway')
    strain(data_collection, ref_omega_dft, ref_omega_pace)
    print('Strain')
    phonos(data_collection)
    print('phonons')
    collect_in_powerpoint(data_collection, RMSE_E, MAE_E, RMSE_F, RMSE_F,pptxname)


######  python {python program name}  {output pptx filename without .pptx} ######
if __name__ == "__main__":
    main(sys.argv[1])